"""PowerPoint converter using LibreOffice bridge with slide screenshot support."""

import logging
import os
import re
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import List, Optional, Tuple

import shutil

from ..core import BaseConverter, ConversionResult, ExtractedImage, LibreOfficeBridge
from ..utils import ImageConverter, OCRClient
from ..utils.ocr_client import is_skip_response, needs_screenshot, get_pptx_prompt

logger = logging.getLogger(__name__)


def _get_parallel_workers() -> int:
    """Get number of parallel workers for Vision API calls."""
    try:
        return int(os.getenv("VISION_API_PARALLEL_WORKERS", "200"))
    except ValueError:
        return 30


def _get_slide_ocr_threshold() -> int:
    """環境変数からスライド単位のOCR閾値を取得.
    
    デフォルト: 99999（実質全スライドOCR）
    全スライドをOCRしてAI判断でスクショ必要性を決定する。
    """
    try:
        return int(os.getenv("SLIDE_OCR_THRESHOLD", "99999"))
    except ValueError:
        return 99999


def _get_complex_slide_threshold() -> int:
    """スクショ対象とするテキスト文字数閾値（デフォルト: 10文字）.
    
    これ以下のテキスト量 = 画像/図中心のスライドと判断してスクショ対象。
    """
    try:
        return int(os.getenv("COMPLEX_SLIDE_TEXT_THRESHOLD", "10"))
    except ValueError:
        return 10


def _get_complex_slide_shape_threshold() -> int:
    """複雑な図スライドと判定するshape数閾値（デフォルト: 5）."""
    try:
        return int(os.getenv("COMPLEX_SLIDE_SHAPE_THRESHOLD", "5"))
    except ValueError:
        return 5


def _is_ocr_all_slides_enabled() -> bool:
    """環境変数から全スライドOCRオプションを取得（デフォルト: false）."""
    return os.getenv("OCR_ALL_SLIDES", "false").lower() in ("true", "1", "yes")


# 図関連キーワード（OCR結果にこれらが含まれていれば複雑な図と判定）
_DIAGRAM_KEYWORDS = [
    "フロー", "フローチャート", "プロセス", "ワークフロー",
    "図", "概念図", "構成図", "組織図", "階層図",
    "チャート", "タイムライン", "ダイアグラム",
    "矢印", "ボックス", "分岐", "処理",
]

def _has_diagram_keywords(text: str) -> bool:
    """テキストに図関連キーワードが含まれているか判定."""
    if not text:
        return False
    text_lower = text.lower()
    return any(kw in text_lower for kw in _DIAGRAM_KEYWORDS)


def _extract_body_text(text: str) -> str:
    """ヘッダー/フッター/テンプレ注意事項を除外して本文テキストを抽出.
    
    除外対象:
    - ページ番号（数字のみの行）
    - 日付パターン（YYYY/MM/DD, YYYY年MM月DD日など）
    - 著作権表記（Copyright, ©など）
    - レベル/カテゴリ表記（Level A, レベルB等の短い行）
    - 「全体目次へ」などのナビゲーションテキスト
    - 空白行
    """
    if not text:
        return ""
    
    lines = text.strip().split("\n")
    body_lines = []
    
    # ヘッダー/フッターパターン
    skip_patterns = [
        r"^\s*\d+\s*$",  # ページ番号のみ
        r"^\s*\d{4}[/\-年]\d{1,2}[/\-月]\d{1,2}日?\s*$",  # 日付
        r"^\s*(Copyright|©|All Rights Reserved)",  # 著作権
        r"^\s*(Level|レベル)\s*[A-Za-z0-9]\s*[:：]?",  # レベル表記
        r"^\s*全体目次へ\s*$",  # ナビゲーション
        r"^\s*WEB\s*機器交換",  # テンプレート固定テキスト
        r"^\s*ソリューション\s*[A-Z]\s*$",  # テンプレート固定テキスト
    ]
    
    compiled_patterns = [re.compile(p, re.IGNORECASE) for p in skip_patterns]
    
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        
        # パターンマッチで除外
        if any(p.match(stripped) for p in compiled_patterns):
            continue
        
        # 非常に短い行（3文字以下）で数字のみの場合は除外
        if len(stripped) <= 3 and stripped.isdigit():
            continue
            
        body_lines.append(stripped)
    
    return " ".join(body_lines)


class PPTConverter(BaseConverter):
    """Convert PowerPoint files to text via LibreOffice PDF conversion.

    LibreOfficeを優先的に使用し、失敗時はpython-pptxにフォールバック。
    複雑な図を含むスライドは自動でスクリーンショットを保存。

    環境変数による設定:
    - SLIDE_OCR_THRESHOLD: スライド単位のOCR閾値（デフォルト: 30文字）
    - OCR_ALL_SLIDES: 全スライドをOCR処理する（デフォルト: false）
    - COMPLEX_SLIDE_TEXT_THRESHOLD: 複雑な図と判定するテキスト閾値（デフォルト: 200文字）
    - COMPLEX_SLIDE_SHAPE_THRESHOLD: 複雑な図と判定するshape数閾値（デフォルト: 10）
    """

    def __init__(self):
        """Initialize PPT converter with LibreOffice bridge."""
        self.bridge = LibreOfficeBridge()

    def convert(self, file_path: Path) -> ConversionResult:
        """Convert PowerPoint file to text with optional slide screenshots.

        Args:
            file_path: Path to PPT/PPTX file

        Returns:
            ConversionResult with text and images
        """
        file_path = Path(file_path).resolve()
        pdf_path: Optional[Path] = None
        log_path: Optional[Path] = None

        # Validate file
        if not file_path.exists():
            return ConversionResult(
                success=False,
                message=f"File not found: {file_path}",
                original_path=file_path,
                file_format="pptx",
            )

        if file_path.suffix.lower() not in {'.ppt', '.pptx', '.odp'}:
            return ConversionResult(
                success=False,
                message=f"Unsupported format: {file_path.suffix}",
                original_path=file_path,
                file_format="pptx",
            )

        try:
            # Convert to PDF using LibreOffice bridge
            pdf_path, success, log_path = self.bridge.convert_to_pdf(file_path)

            if not success:
                # LibreOffice変換失敗時はエラー（python-pptxフォールバックは使用しない）
                logger.error(f"LibreOffice conversion failed for {file_path.name}. See log: {log_path}")
                self._log_tail(log_path)
                return ConversionResult(
                    success=False,
                    message=f"LibreOffice conversion failed for {file_path.name}. Log: {log_path}",
                    original_path=file_path,
                    file_format="pptx",
                )

            # スライド単位でテキスト抽出とOCR判定
            return self._convert_with_slide_level_ocr(pdf_path, file_path, log_path)

        except Exception as e:
            logger.error(f"PPT conversion error for {file_path}: {e}")
            return ConversionResult(
                success=False,
                message=f"Conversion error: {str(e)}",
                original_path=file_path,
                file_format="pptx",
            )
        finally:
            # Clean up temporary PDF if it was auto-generated
            try:
                if pdf_path and pdf_path.exists() and pdf_path.parent.name.startswith('tmp'):
                    pdf_path.unlink()
            except Exception:
                pass  # Ignore cleanup errors

    def _convert_with_slide_level_ocr(
        self, pdf_path: Path, file_path: Path, log_path: Optional[Path]
    ) -> ConversionResult:
        """スライド単位でOCR判定を行い、テキストを抽出する.
        
        複雑な図を含むスライドは自動でスクリーンショットを保存。
        
        Args:
            pdf_path: PDFファイルのパス
            file_path: 元のPPTファイルのパス
            log_path: LibreOfficeのログパス
            
        Returns:
            ConversionResult with text and slide images
        """
        threshold = _get_slide_ocr_threshold()
        ocr_all = _is_ocr_all_slides_enabled()
        complex_text_threshold = _get_complex_slide_threshold()
        complex_shape_threshold = _get_complex_slide_shape_threshold()
        
        # PyMuPDFでページ単位のテキスト抽出
        try:
            import fitz  # PyMuPDF
        except ImportError:
            # PyMuPDFがない場合は従来のpdftotextを使用
            text = self.bridge.extract_text_with_pdftotext(pdf_path)
            if text and len(text.strip()) > 0:
                return ConversionResult(
                    success=True,
                    text=text,
                    message="",
                    original_path=file_path,
                    file_format="pptx",
                )
            return ConversionResult(
                success=False,
                message="PyMuPDF not available and pdftotext failed",
                original_path=file_path,
                file_format="pptx",
            )
        
        try:
            doc = fitz.open(pdf_path)
        except Exception as e:
            logger.error(f"Failed to open PDF: {e}")
            return ConversionResult(
                success=False,
                message=f"Failed to open PDF: {e}",
                original_path=file_path,
                file_format="pptx",
            )
        
        page_count = doc.page_count
        page_texts: List[str] = []
        slides_needing_ocr: List[int] = []  # 0-indexed
        
        # 各ページのテキストを抽出してOCR必要性を判定
        for page_idx in range(page_count):
            page = doc.load_page(page_idx)
            raw_text = page.get_text("text")
            body_text = _extract_body_text(raw_text)
            
            page_texts.append(raw_text)
            
            # OCR判定: 全スライドOCRモードまたは本文が閾値以下
            if ocr_all or len(body_text) <= threshold:
                slides_needing_ocr.append(page_idx)
                logger.debug(f"Slide {page_idx + 1}: needs OCR (body_len={len(body_text)}, threshold={threshold})")
            else:
                logger.debug(f"Slide {page_idx + 1}: text OK (body_len={len(body_text)})")
        
        doc.close()
        
        # OCRが必要なスライドがなければ、通常テキストを返す
        if not slides_needing_ocr:
            combined_text = "\n\n---\n\n".join(
                f"## Slide {idx + 1}\n\n{text}" for idx, text in enumerate(page_texts) if text.strip()
            )
            if combined_text:
                return ConversionResult(
                    success=True,
                    text=combined_text,
                    message="",
                    original_path=file_path,
                    file_format="pptx",
                )
            return ConversionResult(
                success=False,
                message="No text extracted from any slide",
                original_path=file_path,
                file_format="pptx",
            )
        
        logger.info(
            f"OCR needed for {len(slides_needing_ocr)}/{page_count} slides in {file_path.name}"
        )
        
        # OCRが必要なスライドのみ画像化してOCR処理
        image_converter = ImageConverter(self.bridge)
        images, work_dir, img_msg = image_converter.pdf_to_images(pdf_path)
        
        if not images:
            msg = f"OCR fallback skipped: {img_msg}"
            logger.warning(msg)
            self._log_tail(log_path)
            return ConversionResult(
                success=False,
                message=msg,
                original_path=file_path,
                file_format="pptx",
            )
        
        try:
            ocr_client = OCRClient()
        except Exception as e:
            error_msg = f"OCR client init failed: {e}"
            logger.error(error_msg)
            # OCRクライアント初期化失敗時は通常テキストを返す
            combined_text = "\n\n---\n\n".join(
                f"## Slide {idx + 1}\n\n{text}" for idx, text in enumerate(page_texts) if text.strip()
            )
            return ConversionResult(
                success=True,
                text=combined_text,
                message=error_msg,
                original_path=file_path,
                file_format="pptx",
            )
        
        # OCR結果とテキスト抽出結果をマージ（並列処理）
        ocr_results: dict = {}
        max_workers = _get_parallel_workers()
        logger.info(f"Processing {len(slides_needing_ocr)} slides with OCR using {max_workers} parallel workers")
        
        # PPTX専用プロンプトを使用
        pptx_prompt = get_pptx_prompt()
        
        def process_slide_ocr(page_idx: int) -> Tuple[int, Optional[str], Optional[str]]:
            """Worker function to OCR a single slide."""
            if page_idx < len(images):
                img_path = images[page_idx]
                ok, md, err = ocr_client.image_to_markdown(img_path, prompt=pptx_prompt)
                if ok and md:
                    # SKIPレスポンスのチェック
                    if is_skip_response(md):
                        return (page_idx, None, "SKIP")
                    return (page_idx, md, None)
                else:
                    return (page_idx, None, err)
            return (page_idx, None, "Image not found")
        
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {executor.submit(process_slide_ocr, idx): idx for idx in slides_needing_ocr}
            
            completed = 0
            for future in as_completed(futures):
                completed += 1
                try:
                    page_idx, md, err = future.result()
                    if md:
                        ocr_results[page_idx] = md
                        logger.debug(f"Slide {page_idx + 1}: OCR success")
                    elif err == "SKIP":
                        logger.debug(f"Slide {page_idx + 1}: OCR skipped (logo/decoration)")
                    else:
                        logger.warning(f"OCR failed for slide {page_idx + 1}: {err}")
                    
                    if completed % 5 == 0:
                        logger.info(f"OCR progress: {completed}/{len(slides_needing_ocr)} completed")
                except Exception as e:
                    logger.warning(f"Error processing slide: {e}")
        
        logger.info(f"OCR completed: {len(ocr_results)}/{len(slides_needing_ocr)} successful")
        
        # 複雑な図スライドを判定してスクリーンショットを保存
        extracted_images: List[ExtractedImage] = []
        slides_with_screenshot: List[int] = []
        
        for page_idx in slides_needing_ocr:
            body_text = _extract_body_text(page_texts[page_idx])
            ocr_text = ocr_results.get(page_idx, "")
            
            # スクショ判定（AI判断）
            # OCR結果に「SCREENSHOT_NEEDED」が含まれていればスクショ対象
            # AI が「フロー図」「複雑な表」「組織図」等と判断した場合にスクショ
            is_complex = needs_screenshot(ocr_text)
            
            if is_complex and page_idx < len(images):
                img_path = images[page_idx]
                if img_path.exists():
                    try:
                        img_data = img_path.read_bytes()
                        filename = f"slide_{page_idx + 1:03d}.png"
                        extracted_images.append(ExtractedImage(
                            data=img_data,
                            filename=filename,
                            page_or_sheet=f"Slide {page_idx + 1}",
                            position=page_idx + 1,
                            ai_description=ocr_text if ocr_text else None,
                            width=None,
                            height=None,
                            format="png",
                        ))
                        slides_with_screenshot.append(page_idx + 1)
                        logger.debug(f"Slide {page_idx + 1}: screenshot saved (complex diagram)")
                    except Exception as e:
                        logger.warning(f"Failed to save slide {page_idx + 1} screenshot: {e}")
        
        # Clean up images
        if work_dir and work_dir.exists():
            try:
                shutil.rmtree(work_dir)
            except Exception:
                pass
        
        # 最終結果の組み立て
        final_outputs = []
        for page_idx in range(page_count):
            slide_content = f"## Slide {page_idx + 1}\n\n"
            
            # スクリーンショットがある場合は画像参照を追加
            if page_idx + 1 in slides_with_screenshot:
                slide_content += f"![slide_{page_idx + 1:03d}.png](../04_images/slide_{page_idx + 1:03d}.png)\n\n"
            
            if page_idx in ocr_results:
                # OCR結果を使用
                slide_content += ocr_results[page_idx]
            elif page_texts[page_idx].strip():
                # 通常テキストを使用
                slide_content += page_texts[page_idx]
            
            if slide_content.strip():
                final_outputs.append(slide_content)
        
        combined_text = "\n\n---\n\n".join(final_outputs).strip()
        
        if combined_text:
            ocr_count = len(ocr_results)
            msg = f"OCR applied to {ocr_count}/{page_count} slides"
            if slides_with_screenshot:
                msg += f", screenshots saved for {len(slides_with_screenshot)} complex slides"
            return ConversionResult(
                success=True,
                text=combined_text,
                message=msg,
                images=extracted_images,
                original_path=file_path,
                file_format="pptx",
            )
        else:
            warning = f"No content extracted for {file_path.name}"
            logger.warning(warning)
            return ConversionResult(
                success=False,
                message=warning,
                original_path=file_path,
                file_format="pptx",
            )

    @staticmethod
    def _log_tail(log_path: Optional[Path], lines: int = 40) -> None:
        """Log the tail of the LibreOffice log file for quick debugging."""
        try:
            if log_path and log_path.exists():
                content = log_path.read_text(encoding="utf-8", errors="ignore").splitlines()
                tail = "\n".join(content[-lines:])
                logger.debug(f"LibreOffice log tail ({log_path}):\n{tail}")
        except Exception:
            # Logging should never raise in the converter
            pass

    def _convert_with_python_pptx(self, file_path: Path) -> ConversionResult:
        """python-pptxを使用してPPTXを直接処理.

        画像が多いスライドやテキストが少ないスライドは自動でOCR処理します。
        複雑な図を含むスライドはスクリーンショットを保存。

        Args:
            file_path: PPTXファイルのパス

        Returns:
            ConversionResult with text and images
        """
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            return ConversionResult(
                success=False,
                message="python-pptx is not installed. Run: pip install python-pptx",
                original_path=file_path,
                file_format="pptx",
            )

        try:
            prs = Presentation(file_path)
        except Exception as e:
            return ConversionResult(
                success=False,
                message=f"Failed to open PPTX: {e}",
                original_path=file_path,
                file_format="pptx",
            )

        threshold = _get_slide_ocr_threshold()
        ocr_all = _is_ocr_all_slides_enabled()
        complex_text_threshold = _get_complex_slide_threshold()
        complex_shape_threshold = _get_complex_slide_shape_threshold()

        slide_data: List[dict] = []  # {idx, text, images, needs_ocr, shape_count}

        # 各スライドを解析
        for slide_idx, slide in enumerate(prs.slides, 1):
            text_parts = []
            images = []
            shape_count = len(slide.shapes)  # オブジェクト数

            for shape in slide.shapes:
                # テキスト抽出
                if hasattr(shape, "text") and shape.text.strip():
                    body_text = _extract_body_text(shape.text)
                    if body_text:
                        text_parts.append(body_text)

                # 表の処理
                if shape.has_table:
                    table = shape.table
                    table_rows = []
                    for row in table.rows:
                        row_cells = [cell.text.strip() if cell.text else "" for cell in row.cells]
                        table_rows.append(" | ".join(row_cells))
                    if table_rows:
                        text_parts.append("\n".join(table_rows))

                # 画像の収集
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        images.append(image.blob)
                    except Exception:
                        pass

            body_text = " ".join(text_parts)
            needs_ocr = ocr_all or len(body_text) <= threshold

            slide_data.append({
                "idx": slide_idx,
                "text": "\n\n".join(text_parts),
                "images": images,
                "needs_ocr": needs_ocr,
                "shape_count": shape_count,
                "body_text_len": len(body_text),
            })

            if needs_ocr:
                logger.debug(f"Slide {slide_idx}: needs OCR (text_len={len(body_text)}, shapes={shape_count}, images={len(images)})")

        # OCRが必要なスライドを処理
        slides_needing_ocr = [s for s in slide_data if s["needs_ocr"] and s["images"]]

        if slides_needing_ocr:
            logger.info(f"OCR needed for {len(slides_needing_ocr)}/{len(slide_data)} slides")
            self._process_slides_ocr(slides_needing_ocr)

        # 複雑な図スライドを判定
        extracted_images: List[ExtractedImage] = []
        
        for s in slide_data:
            ocr_text = s.get("ocr_result", "")
            body_text = s["text"]
            
            # スクショ判定（AI判断 + shape数）
            # OCR結果に「SCREENSHOT_NEEDED」が含まれていればスクショ対象
            # または shape数15以上なら無条件でスクショ（複雑なオブジェクト）
            is_complex = (
                needs_screenshot(ocr_text) or  # AI判断
                s["shape_count"] >= 15  # shape数15以上
            )
            
            if is_complex and s["images"]:
                # 最初の画像をスライドスクリーンショットとして保存
                img_data = s["images"][0]
                filename = f"slide_{s['idx']:03d}.png"
                extracted_images.append(ExtractedImage(
                    data=img_data,
                    filename=filename,
                    page_or_sheet=f"Slide {s['idx']}",
                    position=s["idx"],
                    ai_description=ocr_text if ocr_text else None,
                    width=None,
                    height=None,
                    format="png",
                ))
                s["has_screenshot"] = True
                logger.debug(f"Slide {s['idx']}: screenshot saved (complex diagram, shapes={s['shape_count']})")

        # 最終結果の組み立て
        final_outputs = []
        for s in slide_data:
            slide_content = f"## Slide {s['idx']}\n\n"
            
            # スクリーンショットがある場合は画像参照を追加
            if s.get("has_screenshot"):
                slide_content += f"![slide_{s['idx']:03d}.png](../04_images/slide_{s['idx']:03d}.png)\n\n"
            
            if s.get("ocr_result"):
                slide_content += s["ocr_result"]
            elif s["text"]:
                slide_content += s["text"]
            
            if slide_content.strip():
                final_outputs.append(slide_content)

        combined_text = "\n\n---\n\n".join(final_outputs).strip()

        if not combined_text:
            return ConversionResult(
                success=False,
                message="No content extracted from PPTX",
                original_path=file_path,
                file_format="pptx",
            )

        ocr_count = sum(1 for s in slide_data if s.get("ocr_result"))
        screenshot_count = sum(1 for s in slide_data if s.get("has_screenshot"))
        msg = f"Converted {len(prs.slides)} slides with python-pptx"
        if ocr_count:
            msg += f", OCR applied to {ocr_count} slides"
        if screenshot_count:
            msg += f", screenshots saved for {screenshot_count} complex slides"

        logger.info(msg)
        return ConversionResult(
            success=True,
            text=combined_text,
            message=msg,
            images=extracted_images,
            original_path=file_path,
            file_format="pptx",
        )

    def _process_slides_ocr(self, slides: List[dict]) -> None:
        """スライド内の画像をOCR処理する.

        Args:
            slides: OCR処理が必要なスライドのリスト（各要素は{"idx", "images", ...}）
        """
        import tempfile

        try:
            ocr_client = OCRClient()
        except Exception as e:
            logger.warning(f"OCR client init failed: {e}")
            return

        max_workers = _get_parallel_workers()

        # 画像をファイルに書き出してOCR処理
        def process_slide(slide: dict) -> Tuple[int, Optional[str]]:
            """1つのスライドの画像をOCR処理."""
            slide_idx = slide["idx"]
            images = slide["images"]

            if not images:
                return slide_idx, None

            ocr_texts = []
            for img_blob in images:
                try:
                    # 一時ファイルに画像を書き出し
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
                        f.write(img_blob)
                        tmp_path = Path(f.name)

                    try:
                        # PPTX専用プロンプトを使用
                        ok, md, _ = ocr_client.image_to_markdown(tmp_path, prompt=get_pptx_prompt())
                        if ok and md:
                            # SKIPレスポンスのチェック
                            if not is_skip_response(md):
                                ocr_texts.append(md)
                    finally:
                        tmp_path.unlink(missing_ok=True)
                except Exception as e:
                    logger.debug(f"Image OCR failed: {e}")

            return slide_idx, "\n\n".join(ocr_texts) if ocr_texts else None

        # 並列処理
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {executor.submit(process_slide, s): s for s in slides}

            for future in as_completed(futures):
                try:
                    slide_idx, ocr_text = future.result()
                    if ocr_text:
                        # 元のslide dictを更新
                        for s in slides:
                            if s["idx"] == slide_idx:
                                s["ocr_result"] = ocr_text
                                break
                        logger.debug(f"Slide {slide_idx}: OCR success")
                except Exception as e:
                    logger.warning(f"Slide OCR error: {e}")
